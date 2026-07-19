# -*- coding: utf-8 -*-
"""
논문 발표 — 연구 배경 슬라이드 v2 (배경답게: 우리 실측 빼고 RQ까지)
법(PIPA §23) → 문제 본질(텍스트형) → 기존 한계(구조적+선행연구) → RQ
우리 실측 52%는 결과 슬라이드로 분리 (배경에서 결과 미리 까지 않음).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY   = RGBColor(0x2B, 0x3C, 0x57)
BRICK  = RGBColor(0xB5, 0x49, 0x4B)
BRICK_BG = RGBColor(0xF7, 0xE7, 0xE7)
INK    = RGBColor(0x26, 0x2B, 0x33)
SUB    = RGBColor(0x5C, 0x66, 0x74)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xFC, 0xFC, 0xFA)
NAVY_BG= RGBColor(0xEC, 0xEF, 0xF4)
GOLD   = RGBColor(0xC8, 0x86, 0x12)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background(); bg.shadow.inherit = False


def _set(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sp = spec[4] if len(spec) > 4 else 1
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
        rPr = r._r.get_or_add_rPr(); rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))


def box(x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, radius=0.08):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set(t.text_frame, lines, align, anchor); return t


def down(cx, y, h=0.22, color=NAVY, w=0.36):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx-w/2), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); s.shadow.inherit = False


# ===== 좌측 챕터 바 =====
box(0, 0, 0.98, 7.5, NAVY)
box(0, 0.30, 0.98, 1.15, BRICK)
tb(0, 0.30, 0.98, 1.15, [("02", 34, True, WHITE)], align=PP_ALIGN.CENTER)
tb(0, 1.55, 0.98, 5.0, [("연\n구\n배\n경", 13, True, RGBColor(0x9B,0xA8,0xBC))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ===== 타이틀 =====
tb(1.22, 0.30, 11.9, 0.86, [("연구 배경 — 왜 «한국어 PII 가드레일»인가", 26, True, NAVY)])
box(1.25, 1.18, 11.83, 0.035, BRICK)

# ===== ① 법 (좌) / ② 텍스트형 본질 (우) =====
sy = 1.42
box(1.25, sy, 5.78, 1.52, NAVY_BG, line=NAVY, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(1.42, sy+0.10, 5.5, 0.40, [("①  법이 «특별 보호»를 요구한다", 13.5, True, NAVY)])
tb(1.42, sy+0.52, 5.5, 0.92, [
    ("개인정보보호법 §23 «민감정보»", 11.5, True, INK, 2),
    ("건강 · 사상/신념 · 정치적 견해 · 성생활 · 범죄경력", 10.5, False, INK, 2),
    ("→ 가장 강하게 보호해야 할 정보로 법에 명시", 10, False, SUB, 0),
], anchor=MSO_ANCHOR.TOP)

box(7.28, sy, 5.78, 1.52, BRICK_BG, line=BRICK, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(7.45, sy+0.10, 5.5, 0.40, [("②  그런데 이 정보들은 «텍스트형»이다", 13.5, True, BRICK)])
tb(7.45, sy+0.52, 5.5, 0.92, [
    ("『계란 알레르기』『양성애』『○○당 당원』", 11.5, True, INK, 2),
    ("→ 정해진 형식·체크섬이 없어 규칙으로 못 잡음", 10.5, False, INK, 2),
    ("(전화·카드 같은 숫자형 PII와 근본적으로 다름)", 10, False, SUB, 0),
], anchor=MSO_ANCHOR.TOP)
down(6.66, sy+1.54, color=NAVY)

# ===== ③ 기존 방어의 구조적 한계 (중앙) — 우리 실측 대신 일반론/선행연구 =====
ey = 3.24
box(1.25, ey, 11.83, 1.92, WHITE, line=NAVY, line_w=1.6, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
tb(1.45, ey+0.10, 11.4, 0.38, [("③  기존 가드레일은 «영어 중심»으로 설계되어 한국어 텍스트형이 사각지대", 14, True, NAVY)])
# 4개 한계 카드
cw = (11.83 - 0.4 - 3*0.16) / 4
cxs = [1.45 + i*(cw+0.16) for i in range(4)]
cards = [
    ("Microsoft Presidio", "한국어 PII 5종만\n(주민·사업자·외국인·면허·여권)\n= 전부 숫자형"),
    ("AWS Bedrock", "정책 모델이 영어로 학습\n→ 한국어 재현율 제한"),
    ("Lakera Guard", "프롬프트 인젝션 전용\n→ PII 탐지 목적 아님"),
    ("선행연구 (KDPII, 2024)", "GPT-4조차 한국어 특화\nPII 33종에 저성능 보고"),
]
for cx, (t, body) in zip(cxs, cards):
    box(cx, ey+0.56, cw, 1.24, NAVY_BG, line=RGBColor(0xC5,0xCE,0xDB), radius=0.10, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lines = [(t, 10.5, True, NAVY, 3)]
    for bl in body.split("\n"):
        lines.append((bl, 9.3, False, INK, 1))
    tb(cx, ey+0.56, cw, 1.24, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
down(6.66, ey+1.94, color=BRICK)

# ===== ④ RQ (하단) =====
ry4 = 5.44
box(1.25, ry4, 11.83, 1.72, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
tb(1.45, ry4+0.12, 11.4, 0.40, [("④  그래서 본 연구는 다음을 묻는다  (연구 질문)", 14, True, RGBColor(0xE6,0xD6,0xA6))])
rq_w = (11.83 - 0.4 - 2*0.16) / 3
rqs = [
    ("RQ1", "기존 다계층 가드레일은\n한국어 텍스트형 PII에\n실제로 얼마나 취약한가?"),
    ("RQ2", "LLM judge(GPT-4o)를\n추가하면 메워지는가?\n비용·지연은 얼마인가?"),
    ("RQ3", "LLM 없이 결정론적으로\n이 공백을 메울 수\n있는가?"),
]
for i, (tag, body) in enumerate(rqs):
    rx = 1.45 + i*(rq_w+0.16)
    box(rx, ry4+0.56, rq_w, 1.04, RGBColor(0x3A,0x4D,0x6B), radius=0.08, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lines = [(tag, 12, True, RGBColor(0xE6,0xD6,0xA6), 3)]
    for bl in body.split("\n"):
        lines.append((bl, 9.8, False, WHITE, 1))
    tb(rx, ry4+0.56, rq_w, 1.04, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = r"C:\litellm\연구배경_v2_배경답게.pptx"
prs.save(out)
print("SAVED:", out)
