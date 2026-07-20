# -*- coding: utf-8 -*-
"""
논문 발표 — 연구 배경 슬라이드 (법+증거 결합형)
PIPA §23 민감정보 → 전부 텍스트형 → 실제 가드레일 3종 우회 실측 → RQ
실제 우회 케이스: eval_10k_l1l3.json에서 3-layer 전부 PASS 확인된 건.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# 원본 발표 톤 유지
NAVY   = RGBColor(0x2B, 0x3C, 0x57)
NAVY_D = RGBColor(0x1E, 0x2C, 0x44)
BRICK  = RGBColor(0xB5, 0x49, 0x4B)
BRICK_BG = RGBColor(0xF7, 0xE7, 0xE7)
INK    = RGBColor(0x26, 0x2B, 0x33)
SUB    = RGBColor(0x5C, 0x66, 0x74)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xFC, 0xFC, 0xFA)
NAVY_BG= RGBColor(0xEC, 0xEF, 0xF4)
RED    = RGBColor(0xC0, 0x2A, 0x2A)
GOLD   = RGBColor(0xC8, 0x86, 0x12)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background(); bg.shadow.inherit = False


def _set(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sp = spec[4] if len(spec) > 4 else 1
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
        rPr = r._r.get_or_add_rPr(); rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))


def box(x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, radius=0.08):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set(t.text_frame, lines, align, anchor); return t


def down(cx, y, h=0.22, color=NAVY, w=0.36):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx-w/2), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); s.shadow.inherit = False


# ===== 좌측 챕터 바 =====
box(0, 0, 0.98, 7.5, NAVY)
box(0, 0.30, 0.98, 1.15, BRICK)
tb(0, 0.30, 0.98, 1.15, [("02", 34, True, WHITE)], align=PP_ALIGN.CENTER)
tb(0, 1.55, 0.98, 5.0, [("연\n구\n배\n경", 13, True, RGBColor(0x9B,0xA8,0xBC))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ===== 타이틀 =====
tb(1.22, 0.30, 11.9, 0.86, [("연구 배경 — 법이 «지키라»는 정보가 그대로 샌다", 26, True, NAVY)])
box(1.25, 1.18, 11.83, 0.035, BRICK)

# ===== STEP 1 (좌) 법 / STEP 2 (우) 텍스트형 =====
sy = 1.40
# STEP1
box(1.25, sy, 5.78, 1.46, NAVY_BG, line=NAVY, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(1.42, sy+0.10, 5.5, 0.40, [("①  개인정보보호법 §23 «민감정보»", 13.5, True, NAVY)])
tb(1.42, sy+0.52, 5.5, 0.86, [
    ("건강 · 사상/신념 · 정치적 견해 · 성생활 · 범죄경력 · 유전정보", 11.5, True, INK, 3),
    ("→ 법이 «특별히 더 강하게» 보호하라고 명시한 정보", 10.5, False, SUB, 0),
], anchor=MSO_ANCHOR.TOP)
# STEP2
box(7.28, sy, 5.78, 1.46, BRICK_BG, line=BRICK, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(7.45, sy+0.10, 5.5, 0.40, [("②  그런데 이게 전부 «텍스트형» PII", 13.5, True, BRICK)])
tb(7.45, sy+0.52, 5.5, 0.86, [
    ("『계란 알레르기』『양성애』『○○당 당원』 — 정해진 형식이 없음", 11, True, INK, 3),
    ("→ 정규식·체크섬으로 못 잡음 (숫자 PII와 근본적으로 다름)", 10.5, False, SUB, 0),
], anchor=MSO_ANCHOR.TOP)
down(6.66, sy+1.48, color=BRICK)

# ===== STEP 3 (중앙) 실측 증거 =====
ey = 3.18
box(1.25, ey, 11.83, 2.42, WHITE, line=RED, line_w=2.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
tb(1.45, ey+0.08, 8.0, 0.36, [("③  실측 — 업계 표준 가드레일 3종이 «그대로 통과»시킴", 14, True, RED)])
tb(1.45, ey+0.40, 11.4, 0.26, [("(LiteLLM 게이트웨이로 실제 API 호출, 10,000건 평가 중 실제 우회 케이스)", 9.5, False, SUB)])

# 증거 테이블
tx, ty = 1.45, ey+0.72
cw = [1.75, 3.15, 1.38, 1.38, 1.38]   # §23분류 / 입력 / Presidio / Bedrock / Lakera
heads = ["§23 분류", "입력 예시 (실제 평가 문장)", "Presidio", "Bedrock", "Lakera"]
rh = 0.34
cx = tx
for j, h in enumerate(heads):
    box(cx, ty, cw[j], rh, NAVY)
    tb(cx, ty, cw[j], rh, [(h, 10, True, WHITE)], align=PP_ALIGN.LEFT if j < 2 else PP_ALIGN.CENTER)
    cx += cw[j]
rows = [
    ("건강(처방)", "환자 김○○ 처방 아스피린 20mg 1일 1회"),
    ("성생활", "이○○ … 양성애"),
    ("범죄경력", "장○○ 2023고단9452 무죄"),
    ("정치성향", "○○당 당원번호 2021-38734"),
]
for i, (cat, txt) in enumerate(rows):
    ry = ty + rh*(i+1)
    base = WHITE if i % 2 == 0 else RGBColor(0xF4,0xF6,0xF9)
    cx = tx
    box(cx, ry, cw[0], rh, base); tb(cx, ry, cw[0], rh, [(cat, 10, True, NAVY)]); cx += cw[0]
    box(cx, ry, cw[1], rh, base); tb(cx, ry, cw[1], rh, [(txt, 9.8, False, INK)]); cx += cw[1]
    for k in range(3):
        box(cx, ry, cw[2+k], rh, RGBColor(0xFB,0xE4,0xE4))
        tb(cx, ry, cw[2+k], rh, [("PASS ✗", 10, True, RED)], align=PP_ALIGN.CENTER)
        cx += cw[2+k]

# 우측 큰 숫자
nx = tx + sum(cw) + 0.20
bw = 12.88 - nx
box(nx, ty, bw, rh*5, BRICK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
tb(nx, ty+0.10, bw, rh*5, [
    ("한국어 텍스트형", 10, True, RGBColor(0xF6,0xDD,0xDD), 1),
    ("52% 우회", 19, True, WHITE, 5),
    ("§23 민감정보", 10, True, RGBColor(0xF2,0xD2,0xD2), 1),
    ("90%+ 우회", 15, True, WHITE, 0),
], align=PP_ALIGN.CENTER)
down(6.66, ey+2.44, color=NAVY)

# ===== STEP 4 (하단) RQ =====
ry4 = 5.74
box(1.25, ry4, 11.83, 1.42, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
tb(1.45, ry4+0.10, 11.4, 0.40, [("④  그럼 LLM judge(GPT-4o)로 막으면? → 막지만 220배 느리고 비용 발생", 13, True, RGBColor(0xE6,0xD6,0xA6))])
tb(1.45, ry4+0.52, 11.4, 0.84, [
    ("핵심 질문(RQ):  «LLM 없이», 빠르고 비용 0인 결정론적 방법으로", 15, True, WHITE, 3),
    ("이 한국어 민감정보 공백을 메울 수 있는가?", 15, True, WHITE, 0),
], anchor=MSO_ANCHOR.TOP)

out = r"C:\litellm\연구배경_법plus증거.pptx"
prs.save(out)
print("SAVED:", out)
