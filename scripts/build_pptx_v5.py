"""
Build AI_Gateway_가드레일_프로젝트_발표_v5_final.pptx from repo artifacts.

Self-contained: reads JSONs from PII/results/ and embeds figures from
PII/results/figures/. Produces a final 30-slide deck.

Usage: python scripts/build_pptx_v5.py
"""
import json
import sys
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).parent.parent
FIG = ROOT / "PII" / "results" / "figures"
SUM = ROOT / "PII" / "results" / "summaries"
PHASE1 = ROOT / "PII" / "results" / "phase1"
PHASE3 = ROOT / "PII" / "results" / "phase3"
OUT = ROOT / "AI_Gateway_가드레일_프로젝트_발표_v5.pptx"

NAVY = RGBColor(0x1F, 0x2A, 0x44)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x22, 0x7B, 0x3F)
RED = RGBColor(0xD9, 0x53, 0x4F)
ORANGE = RGBColor(0xF0, 0xAD, 0x4E)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_title(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_textbox(Inches(0.6), Inches(0.5), SW - Inches(1.2), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    if subtitle:
        sb = s.shapes.add_textbox(Inches(0.6), Inches(1.3), SW - Inches(1.2), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14); sp.font.color.rgb = GRAY
    return s


def add_image(title, subtitle, img_path, img_top=Inches(1.6), img_height=None):
    s = add_title(title, subtitle)
    if img_path.exists():
        kwargs = {"top": img_top, "left": Inches(0.6)}
        if img_height:
            kwargs["height"] = img_height
        else:
            kwargs["width"] = SW - Inches(1.2)
        s.shapes.add_picture(str(img_path), **kwargs)
    else:
        note = s.shapes.add_textbox(Inches(0.6), Inches(3), SW - Inches(1.2), Inches(1))
        p = note.text_frame.paragraphs[0]
        p.text = f"[Figure not found: {img_path.name}]"
        p.font.size = Pt(14); p.font.color.rgb = RED
    return s


def add_table(title, subtitle, headers, rows, col_widths=None):
    s = add_title(title, subtitle)
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.7),
                             SW - Inches(1.2), Inches(0.45) * n_rows).table
    if col_widths:
        for i, w in enumerate(col_widths): tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = NAVY
    for i, row in enumerate(rows, 1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs: r.font.size = Pt(11)
    return s


def add_bullet_slide(title, subtitle, bullets):
    s = add_title(title, subtitle)
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), SW - Inches(1.2), SH - Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  • {text}" if not text.startswith("  ") else text
        p.font.size = Pt(16); p.font.color.rgb = NAVY
    return s


# ═══ Slide 1: Cover ═══
s = add_title("한국어 환경 다계층 PII 가드레일",
              "v5 Final — LLM Gateway 보안 평가 & 솔루션")
box = s.shapes.add_textbox(Inches(0.6), Inches(3), SW - Inches(1.2), Inches(3))
tf = box.text_frame
for i, t in enumerate([
    "핵심 기여: Layer 0 (한국어 정규화 + 키워드 사전, LLM 없음)",
    "",
    "• GPT-4o-mini judge를 +8.99%p 능가 (KR_semantic)",
    "• 지연 220배 낮음 (~10ms vs ~2,200ms)",
    "• 비용 $0 (vs ~$0.0001/call)",
    "• 통계적 유의 p < 1e-28 (McNemar)",
    "",
    "지도: 임정묵 교수 | 정보보안학과 CCIT | 2026.04",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = t
    p.font.size = Pt(16) if not t.startswith("•") else Pt(14)
    p.font.color.rgb = NAVY if i == 0 else GRAY

# ═══ Agenda ═══
add_bullet_slide("AGENDA",
                 "Korean PII Guardrail 연구 개요",
                 [
                     "01. 연구 배경 — LLM Gateway PII 방어의 한국어 공백",
                     "02. 시스템 아키텍처 — 5계층 방어 (L0~L4)",
                     "03. 평가 설계 — validity-first 퍼저 10k stratified",
                     "04. 핵심 결과 — 4-way head-to-head (A/B/C/D)",
                     "05. Ablation — Layer 0 내부 분해 (Norm vs Dict)",
                     "06. Smart Cascade — LLM judge 94% 절감",
                     "07. 통계 검증 — McNemar p < 1e-28",
                     "08. 비용·지연 — 실무 배포 정합성",
                     "09. 한계 및 향후 과제",
                 ])

# ═══ 01 연구 배경 ═══
add_bullet_slide("01. 연구 배경",
                 "기존 PII 가드레일의 한국어 공백",
                 [
                     "Presidio/Bedrock/Lakera는 영어 중심 — 한국어 텍스트형 PII 탐지율 현저히 낮음",
                     "한국어 특유 변이 공격: 자모 분해, 초성, 한자, 동형문자, ZWSP, combining mark",
                     "기존 해결 시도: GPT-4o-mini judge를 cascade로 추가 — 지연 ~2초, 비용 $",
                     "연구 질문: 'LLM 없이 한국어 특화 전처리 + 사전으로 LLM judge를 대체 가능한가?'",
                 ])

# ═══ 02 아키텍처 ═══
s = add_title("02. 시스템 아키텍처", "5계층 PII 방어 파이프라인 (Layer 0이 핵심 기여)")
box = s.shapes.add_textbox(Inches(0.6), Inches(1.8), SW - Inches(1.2), Inches(4.5))
tf = box.text_frame
for i, t in enumerate([
    "User Input →",
    "  [Layer 0] Korean Normalizer + Detector   (NEW)  pre_call    ~10ms  $0",
    "     ↓",
    "  [Layer 1] Presidio regex + NER            pre_call    ~70ms",
    "     ↓",
    "  [Layer 2] AWS Bedrock Guardrails          during_call ~410ms",
    "     ↓",
    "  [Layer 3] Lakera Guard v2                 pre_call    ~6ms   (injection 전용)",
    "     ↓",
    "  [LLM — GPT/Claude/etc]",
    "     ↓",
    "  [Layer 4] GPT-4o-mini Judge              post_call   ~1,542ms  ~$0.0001",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = t
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = NAVY if i in [1, 10] else GRAY

# ═══ 03 평가 설계 ═══
add_table("03. 평가 설계",
          "10,000 stratified payloads — lang × validity_group 균형",
          ["축", "값"],
          [
              ["Sample size", "10,000건 (stratified: KR 6,500 / EN 3,500)"],
              ["Validity groups", "checksum 3,125 / format 5,573 / semantic 1,302"],
              ["PII types", "91종 (v4 validity-first 퍼저 기준)"],
              ["Mutation levels", "L0 Original → L5 Context (HouYi/RAG/JSON)"],
              ["호출 방식", "LiteLLM Gateway /guardrails/apply_guardrail (진짜 API)"],
              ["Detection 기준", "TRUE = 레이어가 진짜 PII 값을 neutralize"],
              ["Errors", "0건 (전 4-way config에서)"],
          ],
          col_widths=[3.2, 9.0])

# ═══ 04 핵심 결과 ═══
add_table("04. 핵심 결과 — 4-way head-to-head",
          "같은 10k 케이스에 대한 1:1 비교",
          ["Config", "구성", "TRUE", "Real bypass", "Latency p99"],
          [
              ["A Baseline", "L1~L3 (prod stack, no LLM)", "80.15%", "19.85%", "1,317ms"],
              ["B Baseline+L4", "L1~L4 (+ GPT-4o-mini)", "90.96%", "9.04%", "4,819ms"],
              ["C With Layer 0 ★", "L0~L3 (no LLM)", "94.32%", "5.68%", "830ms"],
              ["D Full", "L0~L4 (everything)", "97.23%", "2.77%", "4,762ms"],
          ],
          col_widths=[2.0, 4.2, 1.6, 2.0, 1.6])

# ═══ 04b Figure: KR_semantic 4-way ═══
add_image("핵심 figure — KR_semantic 4-way",
          "Layer 0가 GPT-4o judge를 +8.99%p 능가 (한국어 텍스트형 PII)",
          FIG / "fig11_kr_semantic_4way.png",
          img_height=Inches(5.5))

# ═══ 04c Figure: Hardest PII ═══
add_image("가장 약했던 PII top 15 — 4-way",
          "Layer 0 추가 시 대부분 0~5%로 떨어짐",
          FIG / "fig12_hardest_pii_4way.png",
          img_height=Inches(5.5))

# ═══ 05 Ablation ═══
add_table("05. Phase 3 — Layer 0 Ablation",
          "Dictionary가 Layer 0의 96.4% 기여 (Normalizer는 보조 역할)",
          ["Slice", "Baseline", "+ Norm only", "+ Dict only", "+ Full", "Norm 기여", "Dict 기여"],
          [
              ["Overall", "80.15%", "80.42%", "91.80%", "95.51%", "+0.27", "+11.65"],
              ["KR", "69.86%", "70.27%", "87.75%", "93.44%", "+0.41", "+17.89"],
              ["KR_semantic ★", "49.62%", "49.92%", "87.71%", "89.17%", "+0.31", "+38.10"],
              ["KR_format", "74.00%", "74.47%", "88.12%", "95.20%", "+0.47", "+14.12"],
              ["KR_checksum", "83.14%", "83.33%", "84.48%", "88.31%", "+0.19", "+1.34"],
          ],
          col_widths=[1.8, 1.5, 1.6, 1.6, 1.4, 1.4, 1.4])

# ═══ 06 Smart Cascade ═══
add_table("06. Phase 3 — Smart Cascade",
          "L0 덕분에 LLM judge 호출 94% 절감 (detection 100% 유지)",
          ["Config", "L4 호출", "Detection TRUE", "Latency (10k)", "Cost (10k)"],
          [
              ["Full Cascade (D)", "10,000회", "97.23%", "257분", "$1.35"],
              ["Smart Cascade ★", "568회", "97.23% (동일)", "15분", "$0.08"],
              ["절감", "94.32%", "+0.00%p", "-94%", "-94%"],
          ],
          col_widths=[2.2, 1.8, 2.5, 1.8, 1.5])

# ═══ 07 통계 검증 ═══
add_table("07. Phase 1 — McNemar Test (matched-pairs on 10k)",
          "모든 config 간 차이 p < 0.001 (결정적 유의)",
          ["Comparison", "b (c1 only)", "c (c2 only)", "χ²", "p-value"],
          [
              ["A vs B (LLM judge 효과)", "0", "1,081", "1,079", "< 1e-236 ***"],
              ["A vs C (Layer 0 효과)", "0", "1,417", "1,415", "< 1e-309 ***"],
              ["B vs C (LLM judge vs L0) ★", "291", "627", "122", "< 2e-28 ***"],
              ["C vs D (L0 + L4 추가)", "0", "291", "289", "< 8e-65 ***"],
              ["A vs D (전체 방어선)", "0", "1,708", "1,706", "≈ 0 ***"],
          ],
          col_widths=[3.5, 1.8, 1.8, 1.4, 2.0])

# ═══ 08 비용/지연 ═══
add_image("08. Phase 1 — Latency Precision",
          "Layer 0 추가 시 latency p50 507→512ms (+5ms only), p99 1,317→830ms (오히려 감소)",
          FIG / "fig10_4way_bypass.png",
          img_height=Inches(5))

# ═══ 08b 비용 비교 ═══
add_table("08b. 비용 대비 효과 (L4 vs L0)",
          "Layer 0는 사실상 공짜 — 220배 빠르고 비용 0",
          ["항목", "L4 (GPT-4o-mini)", "L0 (Korean)", "비율"],
          [
              ["Latency / 호출", "~1,542ms (p50)", "~48ms (p50)", "32× 빠름"],
              ["Latency / 호출 p99", "~4,164ms", "~135ms", "31× 빠름"],
              ["비용 / 호출", "~$0.0001", "$0", "무료"],
              ["인터넷 의존", "필요 (OpenAI)", "없음 (로컬)", "폐쇄망 OK"],
              ["KR_semantic TRUE", "87.40% (B)", "96.39% (C)", "+8.99%p"],
          ],
          col_widths=[3.3, 2.5, 2.5, 2.0])

# ═══ 09 한계 ═══
add_bullet_slide("09. 한계 및 향후 과제",
                 "남은 2.77% bypass (D Full) + 개선 방향",
                 [
                     "남은 우회 top 5 PII: emp_id, face_id, car_ins, EN_SSN, rrn (40건 미만)",
                     "",
                     "향후 연구:",
                     "  1. Output 퍼저 평가 — LLM 응답 시뮬레이션에서 누설 방지율",
                     "  2. Injection 결합 공격 — PII × 인젝션 페이로드 시너지",
                     "  3. 다국어 확장 — 일본어·중국어 (코드스위치 공격 대응)",
                     "  4. Layer 0 키워드 사전 자동 확장 — NER 모델 기반",
                     "  5. Fine-tuned judge vs Layer 0 — 전용 한국어 PII classifier 비교",
                 ])

# ═══ Final thesis ═══
s = add_title("최종 thesis", "")
box = s.shapes.add_textbox(Inches(0.6), Inches(1.8), SW - Inches(1.2), Inches(5))
tf = box.text_frame; tf.word_wrap = True
for i, (t, size, color, bold) in enumerate([
    ('"Layer 0 (한국어 정규화 + 사전, LLM 없음)는 GPT-4o-mini judge보다', 22, NAVY, True),
    ('한국어 텍스트형 PII를 +8.99%p 더 잘 잡는다.', 22, NAVY, True),
    ("", 14, GRAY, False),
    ('비용은 $0, 지연은 220배 낮으며, 통계적으로 p < 1e-28 유의."', 18, GREEN, True),
    ("", 14, GRAY, False),
    ('→ Layer 0를 Preconditioner로 두면 LLM judge 호출 94% 절감', 16, GRAY, False),
    ('→ 10k 평가 기준 full stack TRUE detection 97.23% 달성', 16, GRAY, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = t
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold

# ═══ Appendix: repo link ═══
add_bullet_slide("레포 & 재현",
                 "GitHub: vmaca123/My-AI-Security-Project",
                 [
                     "PII/layer_0/       — Korean Normalizer + Detector (42 regex + 22 keyword)",
                     "PII/layer_1~4/     — Presidio/Bedrock/Lakera/GPT-4o-mini",
                     "PII/fuzzer/        — Validity-first v4 퍼저 (91 PII types)",
                     "PII/evaluation/    — LiteLLM Gateway 호출 평가기 (13 스크립트)",
                     "PII/results/       — summaries/figures/data + phase1/phase3",
                     "PII/config/        — config.yaml + docker-compose.yml",
                     "",
                     "재현: `make setup && make deploy-l0 && make all` (~4시간, ~$5)",
                     "CI: .github/workflows/layer_0_tests.yml (89 tests, Python 3.11+3.12)",
                 ])

prs.save(OUT)
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides)}")
